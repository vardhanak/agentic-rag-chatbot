import io, uuid
import pandas as pd
import pdfplumber
from pptx import Presentation
from docx import Document

def read_pdf(file_stream: io.BytesIO) -> str:
    text = []
    with pdfplumber.open(file_stream) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text.append(page_text)
    return "\n".join(text)

def read_pptx(file_stream: io.BytesIO) -> str:
    prs = Presentation(file_stream)
    slides_text = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
        slides_text.append("\n".join(parts))
    return "\n\n".join([f"slide {i+1}:\n{t}" for i, t in enumerate(slides_text) if t.strip()])

def read_docx(file_stream: io.BytesIO) -> str:
    doc = Document(file_stream)
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paras)

def read_txt(file_stream: io.BytesIO) -> str:
    return file_stream.read().decode(errors="ignore")

def read_csv(file_stream: io.BytesIO) -> str:
    df = pd.read_csv(file_stream)
    return df.to_csv(index=False)

def infer_and_read(filename: str, file_stream: io.BytesIO):
    ext = filename.split(".")[-1].lower()
    doc_id = f"{uuid.uuid4()}_{filename}"
    if ext == "pdf":
        text = read_pdf(file_stream)
    elif ext == "pptx":
        text = read_pptx(file_stream)
    elif ext == "docx":
        text = read_docx(file_stream)
    elif ext in ("txt","md"):
        text = read_txt(file_stream)
    elif ext == "csv":
        text = read_csv(file_stream)
    else:
        text = read_txt(file_stream)
    return doc_id, text
